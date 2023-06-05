from pptx import Presentation


def get_presentation_as_list_of_slides(path_to_presentation):
    """
    Gets   pptx presentation and read it and return it as list of slides

   Parameters
   ----------
   path_to_presentation : str
        path to pptx presentation

   Returns
   -------
   list
        return the presentation as list of dicts
        example:
        [
           mocks\n
           Python mocks are objects that ' ,

         ' Refactoring code\n
           The process of improving the design ...' ,
           ...
       ]
   """
    presentation = Presentation(path_to_presentation)

    presentation_as_list_of_slides = [get_text_from_slide(slide) for slide in presentation.slides]
    return presentation_as_list_of_slides


def get_text_from_slide(slide):
    """
       Gets  slide and return it as string

      Parameters
      ----------
       slide :  pptx.slide.Slide object

      Returns
      -------
      string
           return the slide as string
            example:

            '  Refactoring code\n
              The process of improving the design ...'
      """
    slide_as_list_of_texts_box = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:

          if not paragraph.text:
                continue

          s = f'{paragraph.text}'.strip()
          slide_as_list_of_texts_box.append(s)

    return '\n'.join(slide_as_list_of_texts_box)




