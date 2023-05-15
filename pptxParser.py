from pptx import Presentation


def get_presentation_as_list_of_slides(path_to_presentation):
    """
    Gets   pptx prestation and read it and retun it as list of slides

   Parameters
   ----------
   path_to_presentation : str
        path to pptx presentation

   Returns
   -------
   list
        return the presentation as list of dicts
        example:
        [
        {'index': 1,
         'text': 'text box:1  mocks\n
                  text box:2  Python mocks are objects that...'}
        {'index': 19,
         'text': 'text box:1  Refactoring code\n
                  text box:2  The process of improving the design ...'}
       ]


   """
    presentation = Presentation(path_to_presentation)

    presentation_as_list_of_slides = [get_text_from_slide(slide) for slide in presentation.slides]

    presentation_as_list_of_slides_with_index = [{'index': index + 1, 'text': slide} for index, slide in
                                                 enumerate(presentation_as_list_of_slides)]

    presentation_as_list_of_slides_filter_empty_slides = filter(lambda s: s['text'] != '',
                                                                presentation_as_list_of_slides_with_index)

    return list(presentation_as_list_of_slides_filter_empty_slides)


def get_text_from_slide(slide):
    """
       Gets  slide and return it as string

      Parameters
      ----------
       slide :  pptx.slide.Slide object

      Returns
      -------
      string
           return the slide as string
            example:

            'text box:1  Refactoring code\n
             text box:2  The process of improving the design ...'
      """
    slide_as_list_of_texts_box = []
    num_slide_run = 1
    for shape in slide.shapes:
        if not shape.has_text_frame:
            return ''

    for paragraph in shape.text_frame.paragraphs:

        if not paragraph.text:
            continue

        s = f'text box:{num_slide_run}  {paragraph.text}'.strip()
        slide_as_list_of_texts_box.append(s)
        num_slide_run += 1

    return '\n'.join(slide_as_list_of_texts_box)


